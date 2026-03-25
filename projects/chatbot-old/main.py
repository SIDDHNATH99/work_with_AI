import json
import logging
import pathlib
import traceback
import uuid
import os
from contextlib import asynccontextmanager
from datetime import datetime, timezone
from typing import Optional, List
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
from starlette.middleware.sessions import SessionMiddleware
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

@asynccontextmanager
async def lifespan(app: FastAPI):
    print("✅ Python server has started!")
    print("🚀 Running on http://localhost:8001")
    yield

app = FastAPI(lifespan=lifespan)

# -------------------- SESSION MIDDLEWARE --------------------
app.add_middleware(SessionMiddleware, secret_key="super-secret-key-change-this")

# -------------------- TEMPLATES --------------------
templates = Jinja2Templates(directory="templates")

# -------------------- CONSTANTS --------------------
MAX_HISTORY = 40
MODEL_NAME  = os.getenv("AZURE_OPENAI_MODEL_NAME", "mistral-small-2503")

# Inference parameters — tune here to change response behaviour
INFERENCE = {
    "temperature":       float(os.getenv("AI_TEMPERATURE",        "0.3")),   # lower = more factual / deterministic
    "top_p":             float(os.getenv("AI_TOP_P",              "0.9")),   # nucleus sampling threshold
    "frequency_penalty": float(os.getenv("AI_FREQUENCY_PENALTY", "0.1")),   # reduce word repetition
    "presence_penalty":  float(os.getenv("AI_PRESENCE_PENALTY",  "0.0")),   # don't force new topics
    "max_tokens":        int(os.getenv("AI_MAX_TOKENS",           "4096")),  # cap per response
}

# System prompt injected at the start of every conversation
SYSTEM_PROMPT = os.getenv("AI_SYSTEM_PROMPT", """\
You are CortexChat, a highly capable AI assistant. Follow these principles:
- Be accurate, concise, and direct. Prefer facts over speculation.
- For code, always include working, complete examples with brief explanations.
- When uncertain, say so clearly rather than guessing.
- Use markdown formatting for structure (headers, bullet points, code blocks).
- Maintain context across the full conversation history provided.
""")

# -------------------- AZURE AI FOUNDRY CLIENT --------------------
client = OpenAI(
    base_url=os.getenv("AZURE_OPENAI_ENDPOINT"),
    api_key=os.getenv("AZURE_OPENAI_API_KEY")
)

# -------------------- PERSISTENT CHAT STORE --------------------
DATA_DIR = pathlib.Path("chats_data")
DATA_DIR.mkdir(exist_ok=True)

def _session_file(session_id: str) -> pathlib.Path:
    # Sanitise: keep only hex chars to prevent path traversal
    safe_id = "".join(c for c in session_id if c in "0123456789abcdef-")
    return DATA_DIR / f"{safe_id}.json"

def load_session_chats(session_id: str) -> dict:
    """Load {chats, active_chat_id} from disk, or return empty defaults."""
    fp = _session_file(session_id)
    if fp.exists():
        try:
            with fp.open("r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            logger.warning(f"Corrupt session file for {session_id[:8]}, resetting.")
    return {"chats": {}, "active_chat_id": None}

def save_session_chats(session_id: str, data: dict) -> None:
    """Persist {chats, active_chat_id} to disk atomically."""
    fp = _session_file(session_id)
    tmp = fp.with_suffix(".tmp")
    with tmp.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False)
    tmp.replace(fp)

# -------------------- MODELS --------------------
class ChatRequest(BaseModel):
    message: str
    chat_id: Optional[str] = None
    image_url: Optional[str] = None   # base64 data URL for vision (validated server-side)

class ChatResponse(BaseModel):
    reply: str
    tokens_used: int
    history_length: int
    chat_id: str
    chat_title: str

class ClearResponse(BaseModel):
    status: str

class ChatListItem(BaseModel):
    id: str
    title: str
    created_at: str
    message_count: int

class ChatsResponse(BaseModel):
    chats: List[ChatListItem]
    active_chat_id: Optional[str]
    active_history: List[dict]

class NewChatResponse(BaseModel):
    chat_id: str

class SwitchChatResponse(BaseModel):
    chat_id: str
    history: List[dict]

# -------------------- HELPERS --------------------
MAX_FILE_BYTES  = 20 * 1024 * 1024   # 20 MB upload limit
MAX_IMAGE_BYTES = 10 * 1024 * 1024   # 10 MB image limit

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".json", ".csv",
    ".html", ".css", ".xml", ".yaml", ".yml", ".sh",
    ".java", ".c", ".cpp", ".h", ".rs", ".go", ".rb",
    ".sql", ".ini", ".toml", ".env",
}

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}
IMAGE_MIME = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".png": "image/png",  ".gif": "image/gif",
    ".webp": "image/webp", ".bmp": "image/bmp",
}

def ensure_session(request: Request) -> str:
    """Ensure a session_id exists in the cookie; return it."""
    if "session_id" not in request.session:
        request.session["session_id"] = str(uuid.uuid4())
        logger.info(f"New session: {request.session['session_id'][:8]}")
    return request.session["session_id"]

def get_store(request: Request) -> dict:
    sid = ensure_session(request)
    store = load_session_chats(sid)
    # Sync active_chat_id from cookie (it's small, kept in session)
    if "active_chat_id" in request.session:
        store["active_chat_id"] = request.session["active_chat_id"]
    return store

def save_store(request: Request, store: dict) -> None:
    sid = request.session["session_id"]
    request.session["active_chat_id"] = store.get("active_chat_id")
    save_session_chats(sid, store)

def create_chat_in_store(store: dict) -> str:
    chat_id = str(uuid.uuid4())
    store["chats"][chat_id] = {
        "title": "New Chat",
        "history": [],
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    store["active_chat_id"] = chat_id
    return chat_id

# -------------------- ROUTES --------------------
# ✅ UI Page
@app.get("/", response_class=HTMLResponse)
async def index(request: Request):
    ensure_session(request)
    return templates.TemplateResponse("index.html", {"request": request})


# ✅ List all chats
@app.get("/chats", response_model=ChatsResponse)
async def list_chats(request: Request):
    store = get_store(request)
    chats = store["chats"]
    active_id = store["active_chat_id"]

    chat_list = [
        ChatListItem(
            id=cid,
            title=chat.get("title", "New Chat"),
            created_at=chat.get("created_at", ""),
            message_count=len(chat.get("history", []))
        )
        for cid, chat in chats.items()
    ]
    chat_list.sort(key=lambda x: x.created_at, reverse=True)

    active_history = []
    if active_id and active_id in chats:
        active_history = chats[active_id].get("history", [])

    return ChatsResponse(chats=chat_list, active_chat_id=active_id, active_history=active_history)


# ✅ Create new chat (defined before /chats/{chat_id}/... to avoid route conflict)
@app.post("/chats/new", response_model=NewChatResponse)
async def new_chat_session(request: Request):
    store = get_store(request)
    chat_id = create_chat_in_store(store)
    save_store(request, store)
    return NewChatResponse(chat_id=chat_id)


# ✅ Switch to a specific chat
@app.post("/chats/{chat_id}/switch", response_model=SwitchChatResponse)
async def switch_chat(request: Request, chat_id: str):
    store = get_store(request)
    if chat_id not in store["chats"]:
        raise HTTPException(status_code=404, detail="Chat not found")
    store["active_chat_id"] = chat_id
    save_store(request, store)
    return SwitchChatResponse(chat_id=chat_id, history=store["chats"][chat_id].get("history", []))


# ✅ Delete a chat
@app.delete("/chats/{chat_id}")
async def delete_chat(request: Request, chat_id: str):
    store = get_store(request)
    if chat_id not in store["chats"]:
        raise HTTPException(status_code=404, detail="Chat not found")
    del store["chats"][chat_id]
    if store["active_chat_id"] == chat_id:
        store["active_chat_id"] = None
    save_store(request, store)
    return {"status": "deleted"}


# ✅ Upload a file — returns extracted text (documents) or base64 data URL (images)
@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    import io, base64

    ext = pathlib.Path(file.filename or "").suffix.lower()
    size_limit = MAX_IMAGE_BYTES if ext in IMAGE_EXTENSIONS else MAX_FILE_BYTES

    raw = await file.read(size_limit + 1)
    if len(raw) > size_limit:
        limit_mb = size_limit // (1024 * 1024)
        raise HTTPException(status_code=413, detail=f"File too large (max {limit_mb} MB for this type)")

    # ---- IMAGES ----
    if ext in IMAGE_EXTENSIONS:
        mime = IMAGE_MIME.get(ext, "image/jpeg")
        b64  = base64.b64encode(raw).decode("ascii")
        return JSONResponse({
            "filename": file.filename,
            "type": "image",
            "image_url": f"data:{mime};base64,{b64}",
            "text": None,
            "chars": len(raw),
        })

    text: str

    # ---- PDF ----
    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(raw))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=422, detail="Could not extract text from PDF")

    # ---- EXCEL ----
    elif ext in (".xlsx", ".xlsm"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    parts.append("\t".join("" if v is None else str(v) for v in row))
            text = "\n".join(parts)
        except Exception:
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=422, detail="Could not extract data from Excel file")

    elif ext == ".xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(file_contents=raw)
            parts = []
            for sheet in wb.sheets():
                parts.append(f"[Sheet: {sheet.name}]")
                for row_idx in range(sheet.nrows):
                    parts.append("\t".join(str(sheet.cell_value(row_idx, c)) for c in range(sheet.ncols)))
            text = "\n".join(parts)
        except Exception:
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=422, detail="Could not extract data from XLS file")

    # ---- WORD ----
    elif ext in (".docx", ".docm"):
        try:
            from docx import Document
            doc = Document(io.BytesIO(raw))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=422, detail="Could not extract text from Word document")

    # ---- POWERPOINT ----
    elif ext in (".pptx", ".pptm"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw))
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
            text = "\n".join(parts)
        except Exception:
            logger.error(traceback.format_exc())
            raise HTTPException(status_code=422, detail="Could not extract text from PowerPoint file")

    # ---- PLAIN TEXT / CODE ----
    elif ext in TEXT_EXTENSIONS or ext == "":
        try:
            text = raw.decode("utf-8", errors="replace")
        except Exception:
            raise HTTPException(status_code=422, detail="Could not read file as text")

    else:
        raise HTTPException(
            status_code=415,
            detail=(
                f"Unsupported file type '{ext}'. "
                "Supported: images (jpg/png/gif/webp/bmp), PDF, Excel (xlsx/xls), "
                "Word (docx), PowerPoint (pptx), and common text/code files."
            )
        )

    return JSONResponse({
        "filename": file.filename,
        "type": "document",
        "image_url": None,
        "text": text,
        "chars": len(text),
    })


def _prepare_history_for_api(history: list) -> list:
    """Strip image data from older messages so API calls stay lean.
    Only the most recent user message retains its full multimodal content.
    """
    result = []
    last_user_idx = max((i for i, m in enumerate(history) if m.get("role") == "user"), default=-1)
    for i, msg in enumerate(history):
        if isinstance(msg.get("content"), list) and i != last_user_idx:
            # Replace image_url parts with a short text note
            text_parts = [p.get("text", "") for p in msg["content"] if p.get("type") == "text"]
            has_image  = any(p.get("type") == "image_url" for p in msg["content"])
            combined   = " ".join(text_parts)
            if has_image:
                combined = (combined + " [image attached]").strip()
            result.append({"role": msg["role"], "content": combined})
        else:
            result.append(msg)
    return result


# ✅ Send message (creates chat automatically if none active)
@app.post("/newchat", response_model=ChatResponse)
async def send_message(request: Request, body: ChatRequest):
    session_id = ensure_session(request)
    user_message = body.message

    # Validate image_url if provided (prevent injection / SSRF)
    image_url = body.image_url
    if image_url is not None:
        if not (image_url.startswith("data:image/") and ";base64," in image_url):
            raise HTTPException(status_code=400, detail="Invalid image_url format")

    store = get_store(request)
    chat_id = body.chat_id or store["active_chat_id"]

    if not chat_id or chat_id not in store["chats"]:
        chat_id = create_chat_in_store(store)
    else:
        store["active_chat_id"] = chat_id

    logger.info(f"[{session_id[:8]}] [{chat_id[:8]}] USER: {user_message[:100]}")

    chat = store["chats"][chat_id]
    history = list(chat.get("history", []))

    # Build user content — multimodal if image is provided
    if image_url:
        user_content: list | str = [
            {"type": "text",      "text": user_message},
            {"type": "image_url", "image_url": {"url": image_url}},
        ]
    else:
        user_content = user_message

    history.append({"role": "user", "content": user_content})

    if len(history) > MAX_HISTORY * 2:
        history = history[-(MAX_HISTORY * 2):]

    api_messages = _prepare_history_for_api(history)

    # Prepend system prompt
    api_messages = [{"role": "system", "content": SYSTEM_PROMPT}] + api_messages

    try:
        response = client.chat.completions.create(
            model=MODEL_NAME,
            messages=api_messages,
            temperature=INFERENCE["temperature"],
            top_p=INFERENCE["top_p"],
            frequency_penalty=INFERENCE["frequency_penalty"],
            presence_penalty=INFERENCE["presence_penalty"],
            max_tokens=INFERENCE["max_tokens"],
        )
        assistant_message = response.choices[0].message.content
        tokens_used = response.usage.total_tokens if response.usage else 0
    except Exception:
        logger.error(traceback.format_exc())
        raise HTTPException(status_code=500, detail="OpenAI API Error")

    history.append({"role": "assistant", "content": assistant_message})

    # Auto-title from first user message
    if chat.get("title") == "New Chat":
        chat["title"] = user_message[:40] + ("..." if len(user_message) > 40 else "")

    chat["history"] = history
    store["chats"][chat_id] = chat
    save_store(request, store)

    logger.info(f"[{session_id[:8]}] AI: {assistant_message[:100]}")

    return ChatResponse(
        reply=assistant_message,
        tokens_used=tokens_used,
        history_length=len(history),
        chat_id=chat_id,
        chat_title=chat["title"]
    )


# ✅ Clear active chat history
@app.post("/clear", response_model=ClearResponse)
async def clear(request: Request):
    store = get_store(request)
    active_id = store["active_chat_id"]
    if active_id and active_id in store["chats"]:
        store["chats"][active_id]["history"] = []
        save_store(request, store)
    return ClearResponse(status="cleared")


# -------------------- RUN --------------------
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="127.0.0.1", port=8001, reload=True)